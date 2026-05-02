import json
import difflib
from pptx import Presentation

# Load DB data
with open('C:/Users/danil/.gemini/antigravity/brain/75a51ba6-67e8-42ec-8eb5-c937290deb2e/.system_generated/steps/133/output.txt', 'r', encoding='utf-8') as f:
    wrapper = json.load(f)
    text = wrapper['result']

start = text.find('[')
end = text.rfind(']') + 1
db_data = json.loads(text[start:end])

# Process PPTX
prs = Presentation('fichas.pptx')
slide_data = []

for slide in prs.slides:
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            # remove weird newlines and extra spaces
            clean_text = ' '.join(shape.text.split())
            texts.append(clean_text)
            
    if not texts: continue
    
    # Heuristic: longest text is the synopsis
    texts.sort(key=len, reverse=True)
    sinopsis = texts[0]
    
    # Title could be any of the others, usually the ones at the beginning or specific ones
    # Let's keep all others as potential titles
    potential_titles = texts[1:]
    slide_data.append({'sinopsis': sinopsis, 'potential_titles': potential_titles})

sql_queries = []
matched_count = 0

for db_item in db_data:
    db_title = db_item['titulo'].lower()
    best_match_sinopsis = None
    best_ratio = 0
    
    for s_data in slide_data:
        for pt in s_data['potential_titles']:
            ratio = difflib.SequenceMatcher(None, db_title, pt.lower()).ratio()
            if ratio > best_ratio:
                best_ratio = ratio
                best_match_sinopsis = s_data['sinopsis']
                
    if best_ratio > 0.8:
        # Escape single quotes
        sinopsis_escaped = best_match_sinopsis.replace("'", "''")
        sql = f"UPDATE public.\\\"FashionFilms\\\" SET sinopsis = '{sinopsis_escaped}' WHERE film_id = '{db_item['film_id']}';"
        sql_queries.append(sql)
        matched_count += 1
    else:
        print(f'No match found for: {db_item["titulo"]} (best ratio: {best_ratio})')

with open('update_pptx.sql', 'w', encoding='utf-8') as f:
    f.write('\n'.join(sql_queries))

print(f'Generated {matched_count} SQL UPDATE statements.')
