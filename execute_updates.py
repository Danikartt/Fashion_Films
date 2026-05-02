import os
import json
import difflib
from supabase import create_client, Client
from pptx import Presentation
from dotenv import load_dotenv

# Load env
load_dotenv('.env.local')
url: str = os.environ.get("NEXT_PUBLIC_SUPABASE_URL")
key: str = os.environ.get("NEXT_PUBLIC_SUPABASE_ANON_KEY")
supabase: Client = create_client(url, key)

# Load DB data
with open('C:/Users/danil/.gemini/antigravity/brain/75a51ba6-67e8-42ec-8eb5-c937290deb2e/.system_generated/steps/133/output.txt', 'r', encoding='utf-8') as f:
    wrapper = json.load(f)
    text = wrapper['result']

start = text.find('[')
end = text.rfind(']') + 1
db_data = json.loads(text[start:end])

# Process PPTX
prs = Presentation('fichas.pptx')
slide_data = []

for slide in prs.slides:
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            clean_text = ' '.join(shape.text.split())
            texts.append(clean_text)
            
    if not texts: continue
    
    texts.sort(key=len, reverse=True)
    sinopsis = texts[0]
    potential_titles = texts[1:]
    slide_data.append({'sinopsis': sinopsis, 'potential_titles': potential_titles})

matched_count = 0

for db_item in db_data:
    db_title = db_item['titulo'].lower()
    best_match_sinopsis = None
    best_ratio = 0
    
    for s_data in slide_data:
        for pt in s_data['potential_titles']:
            ratio = difflib.SequenceMatcher(None, db_title, pt.lower()).ratio()
            if ratio > best_ratio:
                best_ratio = ratio
                best_match_sinopsis = s_data['sinopsis']
                
    if best_ratio > 0.8:
        # Update using supabase client
        response = (
            supabase.table("FashionFilms")
            .update({"sinopsis": best_match_sinopsis})
            .eq("film_id", db_item["film_id"])
            .execute()
        )
        matched_count += 1
        print(f"Updated {matched_count}: {db_item['titulo']}")

print(f'Successfully updated {matched_count} rows in Supabase.')
